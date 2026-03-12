from __future__ import annotations

import json
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


@dataclass(frozen=True)
class SlideSummary:
    title: str
    section: str
    notes: str


@dataclass(frozen=True)
class DeckArtifacts:
    repo_root: Path
    skill_dir: Path
    output_dir: Path
    output_pptx: Path
    output_score: Path


TEAM_MEMBERS: tuple[str, ...] = (
    "Abdul Basit Mohammed",
    "Jude Boachie",
    "Samuel Oduro Duah Boakye",
    "Illiasu Abubakar",
    "Henry Nana Antwi",
)
SECTION_NAMES: tuple[str, ...] = ("Context", "Solution", "Proof", "Retrospective")
PRESENTATION_DATE = "March 12, 2026"


def resolve_artifacts(script_path: Path | None = None) -> DeckArtifacts:
    current_script = (
        Path(__file__).resolve() if script_path is None else script_path.resolve()
    )
    presentation_dir = current_script.parent
    repo_root = presentation_dir.parent.parent
    skill_dir = repo_root / ".agents" / "skills" / "amalitech-pptx-creator"
    output_dir = presentation_dir / "output"
    return DeckArtifacts(
        repo_root=repo_root,
        skill_dir=skill_dir,
        output_dir=output_dir,
        output_pptx=output_dir / "QuickPoll_Team6_Pitch_Deck_v2.pptx",
        output_score=output_dir / "QuickPoll_Team6_Pitch_Deck_v2_score.json",
    )


def build_slide_outline() -> tuple[SlideSummary, ...]:
    return (
        SlideSummary(
            title="QuickPoll helps distributed teams decide faster",
            section="Context",
            notes=(
                "Open with the business value first. Mention that Team 6 built "
                "an internal decision tool for distributed teams and briefly name "
                "every team member once before moving on."
            ),
        ),
        SlideSummary(
            title="Decision friction slows distributed team execution",
            section="Context",
            notes=(
                "Explain why chat threads, calls, and manual follow-up slow teams "
                "down. Keep this non-technical and tie it directly to lost speed "
                "and weak visibility."
            ),
        ),
        SlideSummary(
            title="QuickPoll closes the loop from poll to insight",
            section="Solution",
            notes=(
                "Walk left to right through the user journey. Emphasize that the "
                "product is designed to move teams from question to measurable "
                "decision outcome in one flow."
            ),
        ),
        SlideSummary(
            title="Our MVP covers core polling and analytics needs",
            section="Solution",
            notes=(
                "Use this slide to prove requirements coverage without diving into "
                "implementation detail. Mention the stretch-style extras like "
                "anonymous voting and department targeting as bonus depth."
            ),
        ),
        SlideSummary(
            title="A hybrid demo proves the flow honestly",
            section="Solution",
            notes=(
                "Set expectations clearly: use the Angular UI for login, creation, "
                "and poll listing, then use API or seeded-data evidence for voting, "
                "results, and engagement proof."
            ),
        ),
        SlideSummary(
            title="The architecture keeps QuickPoll reliable",
            section="Proof",
            notes=(
                "Keep this slide at system level. Show how Angular, Spring Boot, "
                "PostgreSQL, and the analytics bootstrap task work together "
                "without overselling unreleased frontend dashboard screens."
            ),
        ),
        SlideSummary(
            title="Trigger-based analytics turns votes into insight",
            section="Proof",
            notes=(
                "Anchor this slide in the current repo truth: create analytics "
                "tables, deploy triggers, run incremental backfill, and let the "
                "database keep metrics current afterward."
            ),
        ),
        SlideSummary(
            title="Quality and delivery practices build trust",
            section="Proof",
            notes=(
                "Present QA and DevOps as one reliability story. Mention test "
                "frameworks, linting, Docker, CI, and deployment readiness in a "
                "tight sequence."
            ),
        ),
        SlideSummary(
            title="Clear ownership kept the team aligned",
            section="Proof",
            notes=(
                "Even with one presenter, this slide should make collaboration "
                "visible. Talk about role coverage, PR flow, migration docs, "
                "standups, and integration checkpoints."
            ),
        ),
        SlideSummary(
            title="QuickPoll maps directly to the evaluator scorecard",
            section="Proof",
            notes=(
                "Read the five dimensions almost exactly as the evaluators will "
                "score them, then point to the evidence already shown in the deck "
                "and the demo flow."
            ),
        ),
        SlideSummary(
            title="Our retrospective shows value, gaps, and next steps",
            section="Retrospective",
            notes=(
                "Be candid here. Mention what went well, what was hard, and that "
                "frontend results and dashboard polish are the main next area to "
                "complete."
            ),
        ),
        SlideSummary(
            title="QuickPoll is ready for broader use and Q&A",
            section="Retrospective",
            notes=(
                "Close with confidence, not with a generic thank-you. Reinforce "
                "that the team delivered a working foundation and is ready for Q&A."
            ),
        ),
    )


def _load_slide_builder(skill_dir: Path) -> Any:
    stub_dir = Path(__file__).resolve().parent / "_pptx_stubs"
    scripts_dir = skill_dir / "scripts"
    template = skill_dir / "assets" / "templates" / "AmaliTech_Blank.pptx"
    brand_config = skill_dir / "assets" / "brand_config.json"
    missing_paths = [
        str(path)
        for path in (stub_dir, scripts_dir, template, brand_config)
        if not path.exists()
    ]
    if missing_paths:
        raise FileNotFoundError(
            "Missing presentation skill asset(s): " + ", ".join(missing_paths)
        )

    sys.path[:0] = [str(stub_dir), str(scripts_dir)]
    module = __import__("pptx_bcg_patterns", fromlist=["BCGSlideBuilder"])
    bcg_slide_builder = module.BCGSlideBuilder

    bcg = bcg_slide_builder(str(template), str(brand_config))
    bcg.preload_icons(
        [
            "chart",
            "check",
            "clock",
            "database",
            "document",
            "gear",
            "globe",
            "lock",
            "person",
            "rocket",
            "security",
            "server",
            "target",
            "team",
        ]
    )
    return bcg


def _team_footer() -> str:
    return "Team 6 | " + " | ".join(TEAM_MEMBERS)


def _configure_console_encoding() -> None:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    if hasattr(sys.stderr, "reconfigure"):
        sys.stderr.reconfigure(encoding="utf-8")


def _notes_text(slide: Any) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def validate_generated_deck(
    pptx_path: Path, outline: tuple[SlideSummary, ...]
) -> dict[str, Any]:
    presentation = Presentation(pptx_path)
    missing_notes = [
        index + 1
        for index, slide in enumerate(presentation.slides)
        if not _notes_text(slide)
    ]
    return {
        "expected_slide_count": len(outline),
        "slide_count": len(presentation.slides),
        "slides_with_notes": len(presentation.slides) - len(missing_notes),
        "notes_complete": not missing_notes,
        "missing_note_slide_numbers": missing_notes,
    }


def build_score_payload(
    outline: tuple[SlideSummary, ...],
    quality_score: dict[str, Any],
    validation: dict[str, Any],
) -> dict[str, Any]:
    return {
        "deck_name": "QuickPoll_Team6_Pitch_Deck_v2",
        "sections": list(SECTION_NAMES),
        "slides": [
            {"title": slide.title, "section": slide.section} for slide in outline
        ],
        "quality": quality_score,
        "validation": validation,
    }


def main() -> None:
    _configure_console_encoding()
    artifacts = resolve_artifacts()
    artifacts.output_dir.mkdir(parents=True, exist_ok=True)
    outline = build_slide_outline()
    bcg = _load_slide_builder(artifacts.skill_dir)
    bcg.set_sections(list(SECTION_NAMES))

    bcg.add_title_slide(
        outline[0].title,
        "A business-first pitch deck for tomorrow's QuickPoll evaluation",
        f"AmaliTech Phase 1 Capstone | {PRESENTATION_DATE}\n{_team_footer()}",
        notes=outline[0].notes,
    )

    bcg.add_scr_slide(
        title=outline[1].title,
        situation=[
            {
                "text": (
                    "Distributed teams lose momentum when decisions are split "
                    "across chats, calls, and scattered follow-ups."
                ),
                "highlights": ["lose momentum"],
            },
            {
                "text": (
                    "Without one auditable vote record, teams struggle to land a "
                    "final answer and confirm who actually participated."
                ),
                "highlights": ["auditable vote record"],
            },
            {
                "text": (
                    "Leads also need immediate engagement visibility instead of "
                    "manual status chasing."
                ),
                "highlights": ["engagement visibility"],
            },
        ],
        resolution=[
            {"type": "heading", "text": "What We Built"},
            {
                "text": (
                    "QuickPoll combines authenticated poll creation, one-vote "
                    "controls, and analytics-ready outputs in one workflow."
                ),
                "highlights": ["one-vote controls"],
            },
            {"type": "heading", "text": "Why It Matters"},
            {
                "text": (
                    "Teams can decide faster, document outcomes clearly, and "
                    "measure participation without extra coordination overhead."
                ),
                "highlights": ["decide faster"],
            },
        ],
        callout=(
            "Fast decisions only help when the result is final, visible, and easy "
            "to trust."
        ),
        source="Project brief, evaluator email, and current Team 6 repository state",
        section=outline[1].section,
        notes=outline[1].notes,
    )

    bcg.add_process_flow_slide(
        title=outline[2].title,
        steps=[
            {
                "title": "Create",
                "icon": "document",
                "items": [
                    "Owners define the question, options, and deadline.",
                    "Single-select or multi-select behavior is configured upfront.",
                ],
            },
            {
                "title": "Target",
                "icon": "team",
                "items": [
                    "Authenticated users access the protected workspace.",
                    "Polls can be company-wide or scoped to departments.",
                ],
            },
            {
                "title": "Vote",
                "icon": "check",
                "items": [
                    "Users submit one final vote within selection rules.",
                    "Expiry and manual close controls protect vote integrity.",
                ],
            },
            {
                "title": "Review",
                "icon": "chart",
                "items": [
                    "Results and participation metrics are exposed for proof.",
                    "Teams move from opinion-heavy debate to visible outcomes.",
                ],
            },
        ],
        source="README.md, PollController, VoteController, METRIC_DEFINITIONS.md",
        takeaway=(
            "QuickPoll links user actions to measurable team decisions in one "
            "end-to-end loop."
        ),
        section=outline[2].section,
        notes=outline[2].notes,
    )

    bcg.add_mece_slide(
        title=outline[3].title,
        boxes=[
            {
                "icon": "security",
                "title": "Auth & Access",
                "bullets": [
                    "JWT registration and login",
                    "Protected user workspace and role-aware access",
                ],
            },
            {
                "icon": "document",
                "title": "Flexible Poll Design",
                "bullets": [
                    "Single-select or multi-select options",
                    "Max selections and expiry included in creation flow",
                ],
            },
            {
                "icon": "team",
                "title": "Audience & Privacy",
                "bullets": [
                    "Department targeting for poll visibility",
                    "Anonymous option available during poll setup",
                ],
            },
            {
                "icon": "chart",
                "title": "Outcome Integrity",
                "bullets": [
                    "Vote finality, close controls, and results API",
                    "Analytics-ready datasets prepared for dashboard reads",
                ],
            },
        ],
        source=(
            "AuthController, PollController, VoteController, create-poll.component.ts"
        ),
        takeaway=(
            "The MVP reaches beyond a basic voting form and delivers a strong "
            "foundation for operational decision-making."
        ),
        section=outline[3].section,
        notes=outline[3].notes,
    )

    bcg.add_process_flow_slide(
        title=outline[4].title,
        steps=[
            {
                "title": "UI Auth",
                "icon": "person",
                "items": [
                    "Open register or login in the Angular interface.",
                    "Enter the protected poll workspace with seeded data ready.",
                ],
            },
            {
                "title": "UI Create",
                "icon": "document",
                "items": [
                    "Create a realistic team decision poll.",
                    "Show options, departments, anonymity, and expiry fields.",
                ],
            },
            {
                "title": "UI List",
                "icon": "globe",
                "items": [
                    "Confirm the poll appears in the list view.",
                    "Use filters to show active or department-relevant polls.",
                ],
            },
            {
                "title": "API & Data Proof",
                "icon": "chart",
                "items": [
                    (
                        "Use Swagger, Postman, or seeded evidence for vote and "
                        "results proof."
                    ),
                    "Tie the response back to analytics tables and engagement metrics.",
                ],
            },
        ],
        source="frontend routes and services, PollController, USER_GUIDE.md",
        takeaway=(
            "The demo stays honest to the current UI while still proving the full "
            "product capability."
        ),
        section=outline[4].section,
        notes=outline[4].notes,
    )

    bcg.add_architecture_slide(
        title=outline[5].title,
        layers=[
            {
                "label": "User Experience",
                "color": "light",
                "components": [
                    {
                        "name": "Angular Frontend",
                        "icon": "globe",
                        "detail": "Register, login, create, list",
                    },
                    {
                        "name": "Auth Guard",
                        "icon": "security",
                        "detail": "Protected workspace access",
                    },
                ],
            },
            {
                "label": "Application API",
                "color": "primary",
                "components": [
                    {
                        "name": "Spring Boot",
                        "icon": "server",
                        "detail": "Auth, polls, votes, results",
                    },
                    {
                        "name": "JWT Security",
                        "icon": "lock",
                        "detail": "Authentication and authorization",
                    },
                ],
            },
            {
                "label": "Operational Data",
                "color": "light",
                "components": [
                    {
                        "name": "PostgreSQL OLTP",
                        "icon": "database",
                        "detail": "Users, polls, options, votes",
                    },
                    {
                        "name": "Poll Scheduler",
                        "icon": "clock",
                        "detail": "Automatic close for expired polls",
                    },
                ],
            },
            {
                "label": "Analytics & Runtime",
                "color": "accent",
                "components": [
                    {
                        "name": "Bootstrap Task",
                        "icon": "gear",
                        "detail": "Tables, triggers, backfill",
                    },
                    {
                        "name": "Analytics Outputs",
                        "icon": "chart",
                        "detail": "Compose locally, staged ECS path",
                    },
                ],
            },
        ],
        connectors=["HTTP + JWT", "JPA/SQL", "Triggers + analytics reads"],
        source="README.md, main.py, USER_GUIDE.md, deploy.yml",
        section=outline[5].section,
        notes=outline[5].notes,
    )

    bcg.add_combo_slide(
        title=outline[6].title,
        left_content=[
            {"type": "heading", "text": "What it maintains"},
            {
                "text": (
                    "The pipeline maintains poll summary, option breakdown, vote "
                    "timeseries, and user participation tables."
                )
            },
            {"type": "heading", "text": "What it measures"},
            {
                "text": (
                    "Key KPIs include participation_rate, vote_percentage, "
                    "votes_in_bucket, total_votes, and unique_voters."
                )
            },
            {"type": "heading", "text": "Why the simplification matters"},
            {
                "text": (
                    "The team moved from an earlier Kafka-oriented idea to a "
                    "trigger-based PostgreSQL design that reduces moving parts "
                    "while keeping analytics current."
                )
            },
            {
                "text": (
                    "Incremental backfill remains as a safety net for historical "
                    "data and reconciliation."
                )
            },
        ],
        right_stats=[
            {"value": "4", "label": "Analytics Tables", "icon": "database"},
            {"value": "5", "label": "Primary KPIs", "icon": "chart"},
            {"value": "1", "label": "Bootstrap Task", "icon": "gear"},
        ],
        source=(
            "main.py, METRIC_DEFINITIONS.md, USER_GUIDE.md, "
            "KAFKA_TO_TRIGGERS_MIGRATION.md"
        ),
        takeaway=(
            "Raw votes become dashboard-ready intelligence through a simpler and "
            "more maintainable analytics path."
        ),
        section=outline[6].section,
        notes=outline[6].notes,
    )

    bcg.add_icon_bullets_slide(
        title=outline[7].title,
        items=[
            {
                "icon": "check",
                "title": "API Test Coverage",
                "description": (
                    "Rest Assured, JUnit 5, and Allure support functional, contract, "
                    "security, and performance validation."
                ),
            },
            {
                "icon": "gear",
                "title": "Python Quality Gates",
                "description": (
                    "Pytest, Ruff, and typed pipeline modules keep the "
                    "data-engineering layer maintainable."
                ),
            },
            {
                "icon": "rocket",
                "title": "CI Discipline",
                "description": (
                    "GitHub Actions runs lint and build checks across backend, "
                    "frontend, and data-engineering services."
                ),
            },
            {
                "icon": "database",
                "title": "Repeatable Environments",
                "description": (
                    "Docker Compose and seeded demo data reduce setup drift before "
                    "presentations and integration testing."
                ),
            },
            {
                "icon": "document",
                "title": "Deployment Readiness",
                "description": (
                    "Infrastructure and handoff guides give the team a clear path "
                    "from local demo to staged deployment."
                ),
            },
        ],
        source="qa/api-tests/README.md, ci.yml, deploy.yml, data-engineering tests",
        section=outline[7].section,
        notes=outline[7].notes,
    )

    bcg.add_mece_slide(
        title=outline[8].title,
        boxes=[
            {
                "icon": "team",
                "title": "Role Coverage",
                "bullets": [
                    (
                        "Frontend, backend, QA, DevOps, and data engineering all "
                        "owned delivery areas"
                    ),
                    (
                        "Cross-functional outputs were integrated instead of "
                        "presented in isolation"
                    ),
                ],
            },
            {
                "icon": "document",
                "title": "PR Flow",
                "bullets": [
                    (
                        "Feature branches and focused pull requests kept reviews "
                        "manageable"
                    ),
                    "Review guidance supported fast, constructive collaboration",
                ],
            },
            {
                "icon": "gear",
                "title": "Shared Documentation",
                "bullets": [
                    "Migration, deployment, and handoff docs reduced ambiguity",
                    "Runbooks supported rehearsal and fallback planning",
                ],
            },
            {
                "icon": "clock",
                "title": "Integration Cadence",
                "bullets": [
                    "Daily standups surfaced blockers early",
                    "Seeded data and checkpoints stabilized the demo story",
                ],
            },
        ],
        source="presentation runbook, PR merge docs, sprint milestones, repo workflow",
        takeaway=(
            "The team worked like one delivery unit, which is why the final story "
            "holds together under time pressure."
        ),
        section=outline[8].section,
        notes=outline[8].notes,
    )

    bcg.add_icon_bullets_slide(
        title=outline[9].title,
        items=[
            {
                "icon": "target",
                "title": "Sprint Goal Achievement",
                "description": (
                    "The team delivered an integrated polling platform with auth, "
                    "creation, voting rules, expiry handling, and analytics outputs."
                ),
            },
            {
                "icon": "globe",
                "title": "User Experience",
                "description": (
                    "The visible UI path is clean for login, creation, and poll "
                    "listing, with a demo designed around what is stable today."
                ),
            },
            {
                "icon": "team",
                "title": "Team Cohesiveness",
                "description": (
                    "Role ownership, reviews, and shared docs made collaboration "
                    "visible in both the product and the presentation."
                ),
            },
            {
                "icon": "rocket",
                "title": "Presentation Delivery",
                "description": (
                    "The demo is preloaded, hybrid by design, and aligned to the "
                    "15-minute presentation window."
                ),
            },
            {
                "icon": "document",
                "title": "Clarity & Structure",
                "description": (
                    "The deck intentionally follows problem, solution, proof, and "
                    "retrospective so evaluators can score it quickly."
                ),
            },
        ],
        source="Evaluator guidance email and this v2 presentation brief",
        section=outline[9].section,
        notes=outline[9].notes,
    )

    bcg.add_combo_slide(
        title=outline[10].title,
        left_content=[
            {"type": "heading", "text": "What went well"},
            {
                "text": (
                    "The team delivered working cross-stack pieces and backed them "
                    "with documentation, automation, and seeded demo data."
                )
            },
            {"type": "heading", "text": "What challenged us"},
            {
                "text": (
                    "Short sprint timing made cross-stack integration and final UI "
                    "polish the hardest part of the delivery window."
                )
            },
            {"type": "heading", "text": "What we would do differently"},
            {
                "text": (
                    "Finish the frontend results and dashboard screens earlier so "
                    "the presentation can rely less on hybrid proof."
                )
            },
            {
                "text": (
                    "Deepen automated end-to-end coverage and keep one shared demo "
                    "environment locked earlier in the sprint."
                )
            },
        ],
        right_stats=[
            {"value": "5", "label": "Delivery Disciplines", "icon": "team"},
            {"value": "12", "label": "Slides in Final Pitch", "icon": "document"},
            {"value": "1", "label": "Hybrid Fallback Plan", "icon": "check"},
        ],
        source="Repository artifacts, presentation notes, and team retrospective",
        takeaway=(
            "The team finished with a credible product foundation and a clear "
            "path to close the remaining presentation-safe gaps."
        ),
        section=outline[10].section,
        notes=outline[10].notes,
    )

    bcg.add_end_slide(outline[11].title, notes=outline[11].notes)

    quality_score = bcg.score_quality()
    bcg.save(str(artifacts.output_pptx))
    validation = validate_generated_deck(artifacts.output_pptx, outline)
    score_payload = build_score_payload(outline, quality_score, validation)
    artifacts.output_score.write_text(
        json.dumps(score_payload, indent=2),
        encoding="utf-8",
    )

    print(f"Presentation saved: {artifacts.output_pptx}")
    print(
        "Quality score: "
        f"{quality_score.get('total')} | passed={quality_score.get('passed')}"
    )
    print(
        "Validation: "
        f"{validation['slide_count']} slides | "
        f"{validation['slides_with_notes']} with notes"
    )
    print(f"Score details saved: {artifacts.output_score}")


if __name__ == "__main__":
    main()
